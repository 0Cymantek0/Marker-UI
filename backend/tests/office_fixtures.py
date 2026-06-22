"""Generate synthetic test fixtures for docx/pptx converters."""

import os
from pathlib import Path
from docx import Document
from docx.shared import Inches

def generate_docx_fixtures(out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1. simple_headings.docx
    doc = Document()
    doc.add_heading("Heading 1", level=1)
    doc.add_paragraph("This is a paragraph under heading 1.")
    doc.add_heading("Heading 2", level=2)
    doc.add_paragraph("This is a paragraph under heading 2.")
    doc.save(out_dir / "simple_headings.docx")

    # 2. tables_links_lists.docx
    doc = Document()
    doc.add_heading("Tables and Lists", level=1)
    doc.add_paragraph("Here is a list:")
    doc.add_paragraph("First item", style="List Bullet")
    doc.add_paragraph("Second item", style="List Bullet")
    
    doc.add_paragraph("Here is a table:")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value A"
    table.cell(1, 1).text = "Value B"
    
    doc.save(out_dir / "tables_links_lists.docx")

    # 3. embedded_text_screenshot.docx
    from PIL import Image
    import io
    
    img = Image.new("RGB", (100, 100), color="blue")
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    
    doc = Document()
    doc.add_heading("Document with Image", level=1)
    doc.add_paragraph("Below is an embedded screenshot image.")
    doc.add_picture(img_byte_arr, width=Inches(2.0))
    doc.save(out_dir / "embedded_text_screenshot.docx")


def generate_pptx_fixtures(out_dir: Path):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from PIL import Image
    import io

    # 1. title_text_notes.pptx
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    slide.shapes.title.text = "Slide Title 1"
    tf = slide.placeholders[1].text_frame
    tf.text = "First paragraph in body text."
    p2 = tf.add_paragraph()
    p2.text = "Second paragraph (bullet point)."
    p2.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Presenter notes for Slide 1."
    prs.save(out_dir / "title_text_notes.pptx")

    # 2. table_chart_images.pptx
    prs = Presentation()
    # Slide 1: Table
    slide1 = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide1.shapes.title.text = "Slide with Table"
    table_shape = slide1.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Cell 1-1"
    table.cell(1, 1).text = "Cell 1-2"

    # Slide 2: Chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Slide with Chart"
    chart_data = CategoryChartData()
    chart_data.categories = ["Category A", "Category B"]
    chart_data.add_series("Series 1", (10, 20))
    chart_data.add_series("Series 2", (15, 25))
    slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(6), Inches(4),
        chart_data
    )

    # Slide 3: Image
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Slide with Image"
    
    img = Image.new("RGB", (100, 100), color="red")
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    slide3.shapes.add_picture(img_byte_arr, Inches(1), Inches(2), width=Inches(3))

    prs.save(out_dir / "table_chart_images.pptx")


if __name__ == "__main__":
    fixtures_dir = Path(__file__).parent / "fixtures" / "conversion"
    generate_docx_fixtures(fixtures_dir)
    generate_pptx_fixtures(fixtures_dir)
    print("Generated DOCX and PPTX fixtures successfully.")
