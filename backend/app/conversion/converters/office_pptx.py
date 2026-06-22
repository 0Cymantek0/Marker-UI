"""OfficePptxConverter — converts PPTX to Markdown using python-pptx.

Traverses slides, extracts titles, text frames (with list formatting), tables,
charts (as markdown tables if data series are present), speaker notes, and
processes embedded images/pictures via EmbeddedImageService.
"""

from __future__ import annotations

import logging
from typing import Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.conversion.registry import BaseConverter
from app.conversion.result import UniversalConversionResult
from app.conversion.embedded_image import EmbeddedImageService
from app.conversion.stream_info import StreamInfo

logger = logging.getLogger(__name__)


class OfficePptxConverter(BaseConverter):
    """Converts PowerPoint presentations (.pptx) to clean Markdown model-free."""

    engine_name = "office_pptx"
    priority = 10
    requires_marker_models = False
    requires_gpu = False

    _EXTENSIONS = frozenset({".pptx"})

    def __init__(self, marker_service: Any = None) -> None:
        self._marker_service = marker_service

    @property
    def supported_extensions(self) -> frozenset[str]:
        return self._EXTENSIONS

    def accepts(self, stream_info: StreamInfo, config: dict[str, Any]) -> bool:
        return stream_info.extension in self._EXTENSIONS

    def convert(
        self,
        filepath: str,
        config: dict[str, Any],
        device: str | None = None,
    ) -> UniversalConversionResult:
        """Parse PPTX, extract text, tables, charts, speaker notes, and process images."""
        try:
            prs = Presentation(filepath)
        except Exception as exc:
            logger.error("Failed to parse pptx file %s: %r", filepath, exc)
            raise

        img_service = EmbeddedImageService(self._marker_service)
        image_counter = 0
        image_metadata = []
        out_images = {}

        def sort_shapes(shapes: list[Any]) -> list[Any]:
            # Sort shapes by visual hierarchy: top-to-bottom, left-to-right.
            # Group shapes into lines by top position within a threshold of 360,000 EMUs (0.4 inches).
            def get_top(s: Any) -> int:
                t = getattr(s, "top", 0)
                return 0 if t is None else t

            def get_left(s: Any) -> int:
                l = getattr(s, "left", 0)
                return 0 if l is None else l

            sorted_by_top = sorted(shapes, key=get_top)
            lines: list[tuple[int, list[Any]]] = []
            for s in sorted_by_top:
                top = get_top(s)
                matched_line = None
                for line in lines:
                    if abs(line[0] - top) < 360000:
                        matched_line = line
                        break
                if matched_line:
                    matched_line[1].append(s)
                else:
                    lines.append((top, [s]))
            
            sorted_shapes = []
            for _, line_shapes in sorted(lines, key=lambda x: x[0]):
                sorted_shapes.extend(sorted(line_shapes, key=get_left))
            return sorted_shapes

        def get_slide_title(slide: Any) -> str:
            title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
            except AttributeError:
                pass
            if not title:
                for shape in slide.shapes:
                    try:
                        if shape.is_placeholder and shape.placeholder_format.type in (1, 3):
                            if hasattr(shape, "text") and shape.text.strip():
                                title = shape.text.strip()
                                break
                    except AttributeError:
                        pass
            return title

        def process_text_frame(text_frame: Any) -> str:
            text_runs = []
            for p in text_frame.paragraphs:
                p_text = p.text.strip()
                if not p_text:
                    continue
                indent = "  " * p.level
                if p.level > 0:
                    text_runs.append(f"{indent}- {p_text}")
                else:
                    text_runs.append(p_text)
            return "\n".join(text_runs)

        def process_table(table: Any) -> str:
            rows_md = []
            for i, row in enumerate(table.rows):
                cells = []
                prev_cell = None
                for cell in row.cells:
                    # python-pptx returns the same _Cell object for every grid
                    # position covered by a merged span; emitting its text for
                    # each column duplicates content. Emit the text once and use
                    # an empty cell for the spanned-over positions.
                    if cell is prev_cell:
                        cells.append("")
                    else:
                        cells.append(cell.text.strip().replace("\n", " "))
                        prev_cell = cell
                rows_md.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
            return "\n".join(rows_md)

        def process_chart(shape: Any) -> str:
            if not getattr(shape, "has_chart", False):
                return ""
            chart = shape.chart
            try:
                categories = []
                if chart.plots and len(chart.plots) > 0 and chart.plots[0].categories:
                    categories = [str(cat.label) if hasattr(cat, "label") else str(cat) for cat in chart.plots[0].categories]
                
                series_data = []
                for s in chart.series:
                    series_data.append((s.name, [val for val in s.values]))
                
                if not categories and not series_data:
                    return ""
                
                headers = ["Series"] + categories
                separator = ["---"] * len(headers)
                rows = []
                rows.append("| " + " | ".join(headers) + " |")
                rows.append("| " + " | ".join(separator) + " |")
                for name, values in series_data:
                    row_cells = [str(name) if name is not None else ""] + [str(val) if val is not None else "" for val in values]
                    rows.append("| " + " | ".join(row_cells) + " |")
                return "\n".join(rows)
            except Exception as e:
                logger.warning("Failed to extract data series from chart shape: %r", e)
                title_text = ""
                try:
                    if chart.has_title and chart.chart_title.text_frame:
                        title_text = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass
                name = f"Chart: {title_text}" if title_text else "Chart"
                return f"\n<!-- [{name}] (Data extraction unavailable) -->\n"

        def extract_shapes(shapes: Any, slide_context: str) -> list[str]:
            nonlocal image_counter
            elements = []
            sorted_s = sort_shapes(list(shapes))
            for shape in sorted_s:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        elements.extend(extract_shapes(shape.shapes, slide_context))
                    except Exception as e:
                        logger.warning("Error processing grouped shape: %r", e)
                    continue

                if getattr(shape, "has_table", False):
                    try:
                        tbl_md = process_table(shape.table)
                        if tbl_md:
                            elements.append(tbl_md)
                    except Exception as e:
                        logger.warning("Error processing table shape: %r", e)
                    continue

                if getattr(shape, "has_chart", False):
                    try:
                        chart_md = process_chart(shape)
                        if chart_md:
                            elements.append(chart_md)
                    except Exception as e:
                        logger.warning("Error processing chart shape: %r", e)
                    continue

                if getattr(shape, "has_text_frame", False):
                    try:
                        txt = process_text_frame(shape.text_frame)
                        if txt:
                            elements.append(txt)
                    except Exception as e:
                        logger.warning("Error processing text shape: %r", e)
                    continue

                is_pic = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE) or hasattr(shape, "image")
                if is_pic:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        content_type = image.content_type
                        ext = "png"
                        if content_type:
                            parts = content_type.split("/")
                            if len(parts) == 2:
                                ext = parts[1]
                        
                        image_counter += 1
                        img_name = f"image_{image_counter}.{ext}"

                        context_text = slide_context
                        if len(elements) > 0:
                            context_text = "\n".join(elements[-2:])

                        try:
                            result = img_service.process_image(
                                image_bytes_or_pil=image_bytes,
                                context_text=context_text,
                                options=config,
                                image_name=img_name,
                            )
                            meta = {
                                "image_name": img_name,
                                "image_type": result.get("image_type", "other"),
                                "confidence": result.get("confidence", 1.0),
                                "model": result.get("model") or ("local-ocr" if result.get("route") == "ocr" else (config.get("vlm_model") or "unknown")),
                                "omitted": result.get("omitted", False),
                                "cost_usd": result.get("cost_usd", 0.0),
                            }
                            if result.get("route") == "skip_decorative" and not result.get("vlm_decided"):
                                meta["model"] = "router-local"
                            image_metadata.append(meta)

                            if not result.get("omitted", False):
                                out_images[img_name] = image_bytes
                                if result.get("markdown"):
                                    elements.append(result["markdown"])
                                else:
                                    elements.append(f"![{img_name}]({img_name})")
                        except Exception as img_exc:
                            logger.warning("Failed to process image %s: %r", img_name, img_exc)
                            out_images[img_name] = image_bytes
                            elements.append(f"![{img_name}]({img_name})")
                    except Exception as e:
                        logger.warning("Error processing picture shape: %r", e)
                    continue

            return elements

        slides_md = []
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            title = get_slide_title(slide)
            if title:
                slides_md.append(f"## Slide {slide_num}: {title}\n")
            else:
                slides_md.append(f"## Slide {slide_num}\n")

            elements = extract_shapes(slide.shapes, title)
            if elements:
                slides_md.append("\n\n".join(elements) + "\n")

            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slides_md.append(f"\n### Notes\n{notes_text}\n")
            except Exception as e:
                logger.warning("Failed to extract notes for slide %d: %r", slide_num, e)

        markdown = "\n".join(slides_md)
        metadata = {}
        if image_metadata:
            metadata["image_understanding"] = image_metadata

        return UniversalConversionResult(
            text=markdown,
            extension="md",
            images=out_images,
            metadata=metadata,
        )
